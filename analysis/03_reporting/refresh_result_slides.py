from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
SLIDES_DIR = ROOT / "results" / "slides"
ASSETS_DIR = SLIDES_DIR / "assets"
SOURCE_PPTX = SLIDES_DIR / "Hedy_area45_inactivation_summary.pptx"
# Note: the current deck has manual PowerPoint edits and is not guaranteed to be
# reproduced exactly by this helper. If you rerun this script, it will overwrite
# the current summary deck with a regenerated version.
OUT_PPTX = SLIDES_DIR / "Hedy_area45_inactivation_summary.pptx"


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slide_ids = list(slide_id_list)
    slide = prs.slides[index]
    rel_id = slide_id_list[index].rId
    prs.part.drop_rel(rel_id)
    slide_id_list.remove(slide_ids[index])


def set_title(slide, text: str) -> None:
    slide.shapes.title.text = text
    tf = slide.shapes.title.text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(24)
            run.font.name = "Aptos"


def set_body(placeholder, lines: list[str], font_size: int = 18) -> None:
    tf = placeholder.text_frame
    tf.clear()
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.name = "Aptos"
    tf.word_wrap = True


def add_title_and_content(prs: Presentation, title: str, lines: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, title)
    set_body(slide.placeholders[1], lines, font_size=20)


def add_picture_fit_box(slide, image_path: Path, left, top, box_w, box_h) -> None:
    with Image.open(image_path) as img:
        img_w, img_h = img.size

    scale = min(box_w / img_w, box_h / img_h)
    pic_w = int(img_w * scale)
    pic_h = int(img_h * scale)
    left = int(left + (box_w - pic_w) / 2)
    top = int(top + (box_h - pic_h) / 2)

    slide.shapes.add_picture(str(image_path), left, top, width=pic_w, height=pic_h)


def add_plot_slide(prs: Presentation, title: str, image_name: str, lines: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, title)

    body = slide.placeholders[1]
    body.left = Inches(0.6)
    body.top = Inches(1.45)
    body.width = Inches(8.2)
    body.height = Inches(1.3)
    set_body(body, lines, font_size=18)

    add_picture_fit_box(slide, ASSETS_DIR / image_name, Inches(0.7), Inches(2.8), Inches(8.6), Inches(4.1))


def add_robustness_text_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(slide, title)

    textbox = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(8.1), Inches(1.5))
    tf = textbox.text_frame
    tf.clear()
    lines = [
        "Quiet-mask analysis: removing smoothed loud epochs from each session did not change the conclusion.",
        "Vet-entry exclusion: removing the session in which a vet entered the room did not change the conclusion.",
    ]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)
        p.font.name = "Aptos"

    rows, cols = 3, 4
    table = slide.shapes.add_table(rows, cols, Inches(0.95), Inches(3.2), Inches(7.2), Inches(2.15)).table
    headers = ["Metric", "Full", "Quiet mask", "Exclude vet"]
    values = [
        ["Net groom duration", "p = 0.0017", "p = 0.0014", "p = 0.0017"],
        ["Grooming reciprocity", "p = 0.0047", "p = 0.0030", "p = 0.0020"],
    ]

    for col, text in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = text
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(18)
                run.font.name = "Aptos"

    for row_idx, row_values in enumerate(values, start=1):
        for col_idx, text in enumerate(row_values):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(18)
                    run.font.name = "Aptos"


def rebuild_result_section(prs: Presentation) -> None:
    while len(prs.slides) > 3:
        delete_slide(prs, 3)

    add_plot_slide(
        prs,
        "DCZ Shifted Grooming Balance Toward Zero",
        "presentation_primary_net_grooming.png",
        [
            "Net groom duration = (received grooming - given grooming) / session duration.",
            "More negative values mean the focal monkey gave more grooming than she received.",
            "0 indicates a balanced exchange.",
        ],
    )

    add_plot_slide(
        prs,
        "DCZ Increased Grooming Reciprocity",
        "presentation_primary_reciprocity.png",
        [
            "Grooming reciprocity = 1 - |received - given| / (received + given).",
            "Higher values indicate a more balanced exchange, independent of direction.",
            "1 indicates perfect reciprocity.",
        ],
    )

    add_robustness_text_slide(
        prs,
        "The Grooming Result Remained Robust in Sensitivity Analyses",
    )

    add_plot_slide(
        prs,
        "A Mild Session-Order Drift Did Not Explain the Grooming Effect",
        "presentation_temporal_dependence.png",
        [
            "We fit separate linear trends within vehicle and DCZ sessions for each primary metric.",
            "The key temporal question was whether either within-condition slope was clearly nonzero.",
            "Neither condition showed a strong slope signal for net grooming or reciprocity.",
            "That makes simple session-order drift a weak explanation for the main grooming effect.",
        ],
    )

    add_plot_slide(
        prs,
        "Total Grooming and Overall Social Engagement Were Less Clearly Affected",
        "presentation_secondary_outcomes.png",
        [
            "Both measures trended lower under DCZ, but neither changed as clearly as the primary grooming metrics.",
            "This makes the main effect look more specific to grooming balance than to total social time.",
        ],
    )

    add_title_and_content(
        prs,
        "Exploratory Follow-Ups Did Not Reveal an Equally Clear Effect",
        [
            "Episode-level grooming turn taking did not show a clear DCZ effect.",
            "Groom-to-nonsocial and nonsocial-to-groom transitions were also not clearly different.",
            "Coarse social-nonsocial transitions were exploratory and less stable than the grooming result.",
            "The clearest effect remained a shift toward more balanced grooming exchange.",
        ],
    )

    add_title_and_content(
        prs,
        "Main Take-Home Message: Area 45 Inactivation Altered Grooming Balance",
        [
            "DCZ shifted grooming from strongly give-biased toward a more balanced exchange.",
            "The two strongest and most robust endpoints were net grooming and grooming reciprocity.",
            "That pattern held across quiet-mask sensitivity, exclusion of the vet-entry session, and the separate within-condition slope checks over session order.",
            "The evidence is strongest for a specific change in grooming balance rather than a broad loss of social interaction.",
        ],
    )


def main() -> None:
    prs = Presentation(SOURCE_PPTX)
    rebuild_result_section(prs)
    prs.save(OUT_PPTX)


if __name__ == "__main__":
    main()

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
FULL_FIGURES = ROOT / "results" / "unblinded" / "full" / "figures"
QUIET_FIGURES = ROOT / "results" / "unblinded" / "quiet_mask" / "figures"
EXCL_FIGURES = ROOT / "results" / "unblinded" / "exclude_vet_entry" / "figures"
FULL_TABLES = ROOT / "results" / "unblinded" / "full" / "tables"
OUT_DIR = ROOT / "results" / "slides"
OUT_PATH = OUT_DIR / "Hedy_area45_inactivation_summary.pptx"


def set_slide_title(slide, title: str) -> None:
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)


def set_body_placeholder(placeholder, bullets: list[str], font_size: int = 18) -> None:
    tf = placeholder.text_frame
    tf.clear()
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = bullet
        p.level = 0
        p.font.size = Pt(font_size)
    tf.word_wrap = True


def add_title_and_content(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, title)
    set_body_placeholder(slide.placeholders[1], bullets)


def add_two_content(prs: Presentation, title: str, left_title: str, left_bullets: list[str], right_title: str, right_bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_slide_title(slide, title)
    set_body_placeholder(slide.placeholders[1], [left_title] + left_bullets, font_size=16)
    set_body_placeholder(slide.placeholders[3], [right_title] + right_bullets, font_size=16)


def fill_picture_placeholder(slide, placeholder_idx: int, image_path: Path) -> None:
    slide.placeholders[placeholder_idx].insert_picture(str(image_path))


def add_picture_with_caption(prs: Presentation, title: str, image_path: Path, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    set_slide_title(slide, title)
    fill_picture_placeholder(slide, 1, image_path)
    set_body_placeholder(slide.placeholders[2], bullets, font_size=16)


def add_comparison_figure_slide(
    prs: Presentation,
    title: str,
    left_header: str,
    left_bullets: list[str],
    left_image: Path,
    right_header: str,
    right_bullets: list[str],
    right_image: Path,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_slide_title(slide, title)
    set_body_placeholder(slide.placeholders[1], [left_header] + left_bullets, font_size=15)
    set_body_placeholder(slide.placeholders[3], [right_header] + right_bullets, font_size=15)

    for ph_idx, image_path in [(2, left_image), (4, right_image)]:
        ph = slide.placeholders[ph_idx]
        slide.shapes.add_picture(str(image_path), ph.left, ph.top, width=ph.width, height=ph.height)


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Area 45 DREADD Inactivation in a Home-Cage Macaque Dyad"
    slide.placeholders[1].text = (
        "Behavioral summary deck\n"
        "Subject: Hedy\n"
        "16 blinded sessions unblinded after locked preprocessing"
    )
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
    for paragraph in slide.placeholders[1].text_frame.paragraphs:
        paragraph.font.size = Pt(18)


def summarize_component_story(path: Path) -> str:
    summary = pd.read_csv(path / "condition_comparison_groom_components.csv")
    give = summary.loc[summary["metric"] == "groom_give_pct_session"].iloc[0]
    receive = summary.loc[summary["metric"] == "groom_receive_pct_session"].iloc[0]
    give_effect = float(give["mean_diff_DCZ_minus_vehicle"])
    receive_effect = float(receive["mean_diff_DCZ_minus_vehicle"])
    give_mag = abs(give_effect)
    receive_mag = abs(receive_effect)

    if give_effect < 0 and receive_effect > 0:
        if give_mag >= 1.5 * max(receive_mag, 1e-9):
            return "The raw decomposition is most consistent with less grooming given under DCZ, plus a smaller increase in grooming received."
        if receive_mag >= 1.5 * max(give_mag, 1e-9):
            return "The raw decomposition is most consistent with more grooming received under DCZ, plus a smaller decrease in grooming given."
        return "The raw decomposition is most consistent with both less grooming given and more grooming received under DCZ."
    if give_effect < 0 and receive_effect <= 0:
        return "The raw decomposition is most consistent with less grooming given under DCZ rather than a clear increase in grooming received."
    if receive_effect > 0 and give_effect >= 0:
        return "The raw decomposition is most consistent with more grooming received under DCZ rather than a clear drop in grooming given."
    return "The raw decomposition does not reduce to a single clean component shift."


def add_take_home(prs: Presentation) -> None:
    component_story = summarize_component_story(FULL_TABLES)
    add_title_and_content(
        prs,
        "Take-Home Points",
        [
            component_story,
            "Composite net grooming and reciprocity still show the clearest session-level condition effects.",
            "The new within-session follow-ups are mechanistic support: they help show how the shift unfolds, but they do not replace the component-first result.",
            "Those effects were robust across quiet-mask sensitivity, exclusion of the vet-entry session, and the separate within-condition slope checks over session order.",
            "Secondary and exploratory analyses still did not reveal an equally clear broad change in total social engagement or coarse transition structure.",
            "The current statistical inference is session-level within one dyad, so it should be framed as strong within-dataset evidence rather than broad population inference.",
        ],
    )


def build_deck() -> None:
    prs = Presentation()
    prs.core_properties.title = "Area 45 DREADD Inactivation Summary"
    prs.core_properties.subject = "Behavioral summary"
    prs.core_properties.author = "OpenAI Codex"

    add_title_slide(prs)

    add_title_and_content(
        prs,
        "Experiment Overview",
        [
            "Bilateral DREADD injections targeted Area 45 in the subject monkey.",
            "Recording began on February 9, 2026, exactly 9 weeks after the December 8, 2025 injections; the subject was then paired daily for 1 hour with her partner for 16 consecutive days.",
            "Half the sessions used intramuscular DCZ 15 minutes before pairing; the other half used vehicle.",
            "Behavior was video-recorded, annotated in BORIS, and analyzed blind to condition before unblinding.",
            "Main question: does Area 45 inactivation alter social exchange, especially grooming reciprocity?",
        ],
    )

    add_two_content(
        prs,
        "Locked Analysis Pipeline",
        "Preprocessing",
        [
            "Analyze only the pairing window, trimmed by 30 s at both ends.",
            "Merge same-label bouts across unlabeled gaps <= 2 s.",
            "Use layered behavior streams: social, activity, attention, atypical.",
            "Bridge short unscored gaps <= 3 s for transition-style analyses.",
        ],
        "Primary endpoints",
        [
            "Net grooming = receive - give, normalized to session duration.",
            "Reciprocity = 1 - abs(receive - give) / (receive + give).",
            "P values from exact two-sided label permutation tests.",
            "Main limitation: repeated sessions from one dyad, so inference is session-level.",
        ],
    )

    add_picture_with_caption(
        prs,
        "Raw Grooming Components",
        FULL_FIGURES / "groom_duration_session_summary.png",
        [
            summarize_component_story(FULL_TABLES),
            "This is now the first interpretive step before discussing composite balance or reciprocity.",
            "Total grooming changed less clearly than the directional give-versus-receive decomposition.",
        ],
    )

    add_picture_with_caption(
        prs,
        "Primary Result",
        FULL_FIGURES / "groom_composite_session_summary.png",
        [
            "After inspecting the raw components, the composite metrics still show a clear shift toward more balanced grooming exchange.",
            "Net grooming shifted strongly toward zero under DCZ.",
            "Reciprocity increased clearly under DCZ.",
            "Full cohort: net grooming p = 0.0017; reciprocity p = 0.0047.",
            "These composite metrics summarize the directional raw change rather than replacing it.",
        ],
    )

    add_picture_with_caption(
        prs,
        "Within-Session Grooming Dynamics",
        FULL_FIGURES / "groom_duration_cumulative_dynamics.png",
        [
            "Condition-averaged cumulative traces show how grooming given, grooming received, and net grooming diverge over the course of a session.",
            "This panel is descriptive and is meant to support mechanism, not replace the session-level condition comparison.",
            "It helps ask whether the raw component shift looks immediate, progressive, or partly feedback-driven within sessions.",
        ],
    )

    add_comparison_figure_slide(
        prs,
        "Robustness Checks",
        "Quiet-mask sensitivity",
        [
            "Removing smoothed loud epochs from each session did not change the primary picture.",
            "Quiet-masked full cohort: net grooming p = 0.0014; reciprocity p = 0.0030.",
        ],
        QUIET_FIGURES / "groom_composite_session_summary.png",
        "Exclude vet-entry session",
        [
            "Dropping session 596273 left the primary effect intact or slightly stronger.",
            "Exclude-vet cohort: net grooming p = 0.0017; reciprocity p = 0.0020.",
        ],
        EXCL_FIGURES / "groom_composite_session_summary.png",
    )

    add_picture_with_caption(
        prs,
        "Temporal Dependence",
        FULL_FIGURES / "temporal_dependence.png",
        [
            "The temporal check now asks whether vehicle and DCZ each show a significant within-condition slope over session order.",
            "Full cohort: neither vehicle nor DCZ slope is clearly significant for net grooming or reciprocity.",
            "That makes simple session-order drift a weak explanation for the main grooming result.",
            "Exclude-vet cohort shows the same pattern.",
        ],
    )

    add_title_and_content(
        prs,
        "Exploratory Follow-Ups",
        [
            "The directional grooming follow-up asks whether Hedy-start and Hooke-start episodes differ under DCZ.",
            "That analysis is meant as mechanistic support for the component-first result rather than a replacement for it.",
            "The older episode-level turn-taking and grooming-to-nonsocial transition summaries remain available as broader exploratory context.",
            "Coarse social/nonsocial macro-transitions were not convincing in the full set; they looked somewhat stronger after excluding the vet-entry session, but remain exploratory.",
            "Attention-to-outside behavior was strongly influenced by the vet-entry session, consistent with the video annotation.",
        ],
    )

    add_picture_with_caption(
        prs,
        "Supporting Visuals",
        FULL_FIGURES / "groom_directional_followup.png",
        [
            "This panel summarizes the directional grooming follow-up metrics.",
            "These measures ask whether same-episode reciprocation depends on whether grooming began with give or receive.",
            "They are useful mechanistic context, but the strongest result remains the raw-to-composite session-level grooming shift.",
        ],
    )

    add_take_home(prs)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name is None:
                        run.font.name = "Aptos"

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)


if __name__ == "__main__":
    build_deck()
